"""
Document processors for extracting text and metadata from various file types.
"""
import os
import io
import re
import json
import math
import chardet
from typing import Tuple, BinaryIO, Dict, Any
from datetime import datetime

import PyPDF2
import docx
import markdown
import csv
import pandas as pd
from bs4 import BeautifulSoup

from app.schemas.document import DocumentType, DocumentMetadata
from app.services.documents.base import logger

def process_text_metadata(content: str, existing_metadata=None) -> DocumentMetadata:
    """Process text content to extract metadata"""
    # Convert existing_metadata to DocumentMetadata if it's a dict
    if isinstance(existing_metadata, dict):
        existing_metadata = DocumentMetadata(**existing_metadata)
    elif existing_metadata is None:
        existing_metadata = DocumentMetadata()
        
    try:
        # Count words
        word_count = len(content.split())
        existing_metadata.word_count = word_count
        
        # Estimate reading time
        existing_metadata.read_time_minutes = math.ceil(word_count / 200)
        
        return existing_metadata
    except Exception as e:
        logger.error(f"Error processing text metadata: {e}")
        return existing_metadata

def get_document_type_from_extension(extension: str) -> DocumentType:
    """Map file extension to document type"""
    extension_map = {
        "txt": DocumentType.TEXT,
        "pdf": DocumentType.PDF,
        "docx": DocumentType.DOCX,
        "doc": DocumentType.DOCX,
        "md": DocumentType.MARKDOWN,
        "csv": DocumentType.CSV,
        "json": DocumentType.JSON,
        "xls": DocumentType.EXCEL,
        "xlsx": DocumentType.EXCEL,
        "html": DocumentType.HTML,
        "htm": DocumentType.HTML,
    }
    
    return extension_map.get(extension.lower(), DocumentType.OTHER)

def extract_text_from_file(file: BinaryIO, filename: str) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from a file based on its type"""
    # Get file extension and determine document type
    extension = os.path.splitext(filename)[1].lower()
    file_size = file.seek(0, os.SEEK_END)
    file.seek(0)  # Reset file pointer after getting size
    
    metadata = DocumentMetadata(
        file_size=file_size,
        title=os.path.splitext(filename)[0],  # Default title is filename without extension
        file_type=extension
    )
    
    try:
        # Extract text based on file type
        if extension == '.txt':
            # For plain text files
            # Try multiple encodings to handle different character sets
            raw_data = file.read()
            detected = chardet.detect(raw_data)
            encoding = detected['encoding'] or 'utf-8'
            
            try:
                content = raw_data.decode(encoding)
            except UnicodeDecodeError:
                # Fallback to latin-1 which can decode any byte sequence
                content = raw_data.decode('latin-1')
                
            metadata = process_text_metadata(content, metadata)
            return content, metadata
            
        elif extension == '.pdf':
            return _extract_from_pdf(file, metadata)
                
        elif extension in ['.docx', '.doc']:
            return _extract_from_docx(file, metadata)
                
        elif extension == '.md':
            return _extract_from_markdown(file, metadata)

        elif extension == '.csv':
            return _extract_from_csv(file, metadata)

        elif extension in ['.xls', '.xlsx']:
            return _extract_from_excel(file, metadata)

        elif extension in ['.json']:
            return _extract_from_json(file, metadata)

        elif extension in ['.html', '.htm']:
            return _extract_from_html(file, metadata)
        
        elif extension in ['.xml']:
            return _extract_from_xml(file, metadata)
            
        elif extension in ['.pptx', '.ppt']:
            return _extract_from_ppt(file, metadata)
        
        elif extension in ['.epub', '.mobi']:
            return _extract_from_ebook(file, metadata)
            
        else:
            logger.warning(f"Unsupported file type: {extension}")
            return f"Unsupported file type: {extension}. The application currently supports txt, pdf, docx, md, csv, excel, json, html, xml, powerpoint, and ebook formats.", metadata
    
    except Exception as e:
        error_msg = f"Error extracting text from file: {str(e)}"
        logger.error(error_msg)
        return error_msg, metadata

def _extract_from_pdf(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from PDF files"""
    try:
        file_bytes = io.BytesIO(file.read())
        pdf_reader = PyPDF2.PdfReader(file_bytes)
        
        # Extract metadata
        info = pdf_reader.metadata
        if info:
            metadata.title = info.title or metadata.title
            metadata.author = info.author
            metadata.created_at = info.creation_date
            metadata.modified_at = info.modification_date
        
        metadata.page_count = len(pdf_reader.pages)
        
        # Extract text
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text() + "\n"
        
        # Process text metadata
        metadata = process_text_metadata(text, metadata)
        return text, metadata
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return f"Error extracting text: {str(e)}", metadata

def _extract_from_docx(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from Word documents"""
    try:
        file_bytes = io.BytesIO(file.read())
        doc = docx.Document(file_bytes)
        
        # Extract metadata
        core_properties = doc.core_properties
        metadata.title = core_properties.title or metadata.title
        metadata.author = core_properties.author
        metadata.created_at = core_properties.created
        metadata.modified_at = core_properties.modified
        metadata.description = core_properties.comments
        
        # Extract text
        content = "\n".join([para.text for para in doc.paragraphs])
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from Word document: {e}")
        return f"Error extracting text: {str(e)}", metadata

def _extract_from_markdown(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from Markdown files"""
    try:
        md_text = file.read().decode('utf-8')
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, features="html.parser")
        
        # Extract potential title from first heading
        headings = soup.find_all(['h1', 'h2'])
        if headings and not metadata.title:
            metadata.title = headings[0].get_text()
        
        # Extract text
        content = soup.get_text()
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from Markdown: {e}")
        return f"Error extracting text: {str(e)}", metadata

def _extract_from_csv(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from CSV files"""
    try:
        # Read the file content
        raw_data = file.read()
        
        # Try multiple encodings to handle different character sets
        encodings_to_try = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        content = None
        
        # Try each encoding until one works
        for encoding in encodings_to_try:
            try:
                text_io = io.StringIO(raw_data.decode(encoding))
                # Test if we can read it as CSV
                reader = csv.reader(text_io)
                rows = list(reader)
                
                if rows:  # If we successfully parsed rows, use this encoding
                    text_io.seek(0)  # Reset the StringIO object
                    reader = csv.reader(text_io)
                    rows = list(reader)
                    content = rows
                    break
            except (UnicodeDecodeError, csv.Error):
                continue
        
        # If all encodings failed, use a fallback approach
        if content is None:
            logger.warning(f"Could not parse CSV with standard encodings, using fallback method")
            try:
                # Try to detect encoding
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'latin-1'  # Use latin-1 as last resort
                
                # Try to decode with detected encoding
                text = raw_data.decode(encoding, errors='replace')
                
                # Simple CSV parsing by splitting on newlines and commas
                lines = text.split('\n')
                content = [line.split(',') for line in lines if line.strip()]
            except Exception as e:
                logger.error(f"Fallback CSV parsing failed: {e}")
                content = [["Error parsing CSV file"]]
        
        # Extract headers and first few rows as a sample
        headers = content[0] if content else []
        sample_rows = content[1:6] if len(content) > 1 else []
        
        # Format as text
        text_content = f"Headers: {', '.join(headers)}\n\n"
        text_content += "Sample data:\n"
        for i, row in enumerate(sample_rows):
            text_content += f"Row {i+1}: {', '.join(row)}\n"
            
        # Add summary
        text_content += f"\nTotal rows: {len(content) - 1}"
        
        # Process text metadata
        metadata = process_text_metadata(text_content, metadata)
        metadata.description = f"CSV file with {len(content) - 1} rows and {len(headers)} columns"
        return text_content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from CSV: {e}")
        # Return a more informative error message
        error_msg = f"Error extracting text from CSV: {str(e)}. The file may be corrupted or in an unsupported format."
        return error_msg, metadata

def _extract_from_excel(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from Excel files"""
    try:
        file_bytes = io.BytesIO(file.read())
        df = pd.read_excel(file_bytes)
        
        # Extract basic info
        content = f"Worksheet with {len(df)} rows and {len(df.columns)} columns.\n\n"
        content += f"Column names: {', '.join(df.columns)}\n\n"
        
        # Add a sample of the data (first 5 rows)
        content += "Sample data:\n"
        content += df.head(5).to_string()
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        metadata.description = f"Excel file with {len(df)} rows and {len(df.columns)} columns"
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {e}")
        return f"Error extracting text: {str(e)}", metadata

def _extract_from_json(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from JSON files"""
    try:
        json_text = file.read().decode('utf-8')
        data = json.loads(json_text)
        
        # Pretty print the JSON for better readability
        formatted_json = json.dumps(data, indent=2)
        
        # Limit the content to a reasonable size if too large
        if len(formatted_json) > 10000:
            content = formatted_json[:10000] + "\n...(truncated)..."
        else:
            content = formatted_json
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        
        # Try to determine if it's an array or object
        if isinstance(data, list):
            metadata.description = f"JSON array with {len(data)} items"
        else:
            metadata.description = f"JSON object with {len(data.keys())} keys"
        
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from JSON: {e}")
        return f"Error extracting text: {str(e)}", metadata

def _extract_from_html(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from HTML files"""
    try:
        html_text = file.read().decode('utf-8')
        soup = BeautifulSoup(html_text, 'html.parser')
        
        # Extract title
        title_tag = soup.find('title')
        if title_tag:
            metadata.title = title_tag.get_text()
        
        # Extract meta description
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        if meta_desc and 'content' in meta_desc.attrs:
            metadata.description = meta_desc['content']
        
        # Extract text content (remove scripts, styles, etc.)
        for script in soup(["script", "style"]):
            script.extract()
        
        # Get text and clean it up
        content = soup.get_text(separator='\n')
        content = re.sub(r'\n+', '\n', content).strip()
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from HTML: {e}")
        return f"Error extracting text: {str(e)}", metadata 

def _extract_from_xml(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from XML files"""
    try:
        xml_text = file.read().decode('utf-8')
        soup = BeautifulSoup(xml_text, features="xml")
        
        # Extract text
        content = soup.get_text()
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        return content, metadata
    except Exception as e:
        logger.error(f"Error extracting text from XML: {e}")
        return f"Error extracting text from XML: {str(e)}", metadata

def _extract_from_ppt(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from PowerPoint files"""
    try:
        # Import here to avoid dependency issues if python-pptx is not installed
        from pptx import Presentation
        
        file_bytes = io.BytesIO(file.read())
        presentation = Presentation(file_bytes)
        
        # Extract text from slides
        content = []
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
        
        text = "\n\n".join(content)
        
        # Update metadata
        metadata.page_count = len(presentation.slides)
        metadata = process_text_metadata(text, metadata)
        
        return text, metadata
    except ImportError as e:
        logger.error(f"python-pptx library not installed: {e}")
        return "PowerPoint extraction requires python-pptx library. Please install it with 'pip install python-pptx'.", metadata
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint: {e}")
        return f"Error extracting text from PowerPoint: {str(e)}", metadata
        
def _extract_from_ebook(file: BinaryIO, metadata: DocumentMetadata) -> Tuple[str, DocumentMetadata]:
    """Extract text and metadata from eBook files (epub, mobi)"""
    try:
        # Import here to avoid dependency issues if ebooklib is not installed
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        
        file_bytes = io.BytesIO(file.read())
        
        # Parse EPUB
        book = epub.read_epub(file_bytes)
        
        # Extract metadata
        metadata.title = book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else metadata.title
        if book.get_metadata('DC', 'creator'):
            metadata.author = book.get_metadata('DC', 'creator')[0][0]
        
        # Extract content from HTML documents
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                chapters.append(soup.get_text())
        
        content = "\n\n".join(chapters)
        
        # Process text metadata
        metadata = process_text_metadata(content, metadata)
        
        return content, metadata
    except ImportError as e:
        logger.error(f"ebooklib library not installed: {e}")
        return "eBook extraction requires ebooklib and BeautifulSoup4 libraries. Please install with 'pip install ebooklib beautifulsoup4'.", metadata
    except Exception as e:
        logger.error(f"Error extracting text from eBook: {e}")
        return f"Error extracting text from eBook: {str(e)}", metadata 